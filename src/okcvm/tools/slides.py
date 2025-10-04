"""Convert Tailwind-flavoured HTML into a PPTX deck."""

from __future__ import annotations

import datetime as _dt
from pathlib import Path
from typing import List, Tuple

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt

from .base import Tool, ToolError, ToolResult


def _parse_slides(html: str) -> List[BeautifulSoup]:
    soup = BeautifulSoup(html, "html.parser")
    slides = soup.select(".ppt-slide")
    if not slides:
        raise ToolError("No elements with class 'ppt-slide' were found in the HTML")
    return slides


def _extract_slide_content(
    slide_markup: BeautifulSoup, index: int
) -> Tuple[str, List[str], List[str]]:
    title_tag = slide_markup.find(["h1", "h2", "h3"])
    title = title_tag.get_text(strip=True) if title_tag else ""
    if not title:
        title = f"Slide {index + 1}"

    paragraphs = [text for p in slide_markup.find_all("p") if (text := p.get_text(strip=True))]
    list_items = [text for li in slide_markup.find_all("li") if (text := li.get_text(strip=True))]

    return title, paragraphs, list_items


def _default_output_path() -> Path:
    timestamp = _dt.datetime.utcnow().strftime("%Y%m%d-%H%M%S")
    output_dir = Path.cwd() / "generated_slides"
    output_dir.mkdir(parents=True, exist_ok=True)
    return output_dir / f"slides-{timestamp}.pptx"


def _add_textbox(slide, text: str, left: float, top: float, width: float, height: float, font_size: int = 32) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.text = text
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)


class SlidesGeneratorTool(Tool):
    name = "mshtools-slides_generator"

    def call(self, **kwargs) -> ToolResult:  # type: ignore[override]
        html = kwargs.get("html") or kwargs.get("content")
        output_path = kwargs.get("output_path")
        if not html:
            raise ToolError("'html' is required")

        slides = _parse_slides(str(html))
        presentation = Presentation()
        blank_layout = presentation.slide_layouts[6]

        preview_slides = []

        for slide_index, slide_markup in enumerate(slides):
            slide = presentation.slides.add_slide(blank_layout)
            title, paragraphs, list_items = _extract_slide_content(slide_markup, slide_index)

            if title:
                _add_textbox(slide, title, left=0.5, top=0.3, width=9.0, height=1.2, font_size=40)

            for idx, text in enumerate(paragraphs):
                _add_textbox(
                    slide,
                    text,
                    left=0.8,
                    top=1.8 + 0.8 * idx,
                    width=8.5,
                    height=0.7,
                    font_size=24,
                )

            for li_index, bullet in enumerate(list_items):
                _add_textbox(
                    slide,
                    f"• {bullet}",
                    left=1.0,
                    top=2.5 + li_index * 0.6,
                    width=8.0,
                    height=0.6,
                    font_size=22,
                )

            outline = paragraphs + list_items
            preview_slides.append({"title": title, "bullets": outline})

        path = Path(output_path).expanduser() if output_path else _default_output_path()
        presentation.save(path)
        return ToolResult(
            success=True,
            output=f"Slides saved to {path}",
            data={"path": str(path), "slides": preview_slides},
        )


__all__ = ["SlidesGeneratorTool"]

